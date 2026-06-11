from pathlib import Path
import re
from pypdf import PdfReader
from docx import Document
from pptx import Presentation


def parse_txt(path: str) -> str:
    return Path(path).read_text(encoding="utf-8", errors="ignore")


def parse_markdown(path: str) -> str:
    return parse_txt(path)


def parse_pdf(path: str) -> str:
    return "\n".join(page.extract_text() or "" for page in PdfReader(path).pages)


def parse_docx(path: str) -> str:
    return "\n".join(p.text for p in Document(path).paragraphs)


def parse_pptx(path: str) -> str:
    return "\n".join(shape.text for slide in Presentation(path).slides for shape in slide.shapes if hasattr(shape, "text"))


def clean_text(text: str) -> str:
    return re.sub(r"\n{3,}", "\n\n", re.sub(r"[ \t]+", " ", text)).strip()


def extract_text(file_path: str) -> str:
    suffix = Path(file_path).suffix.lower()
    parsers = {".txt": parse_txt, ".md": parse_markdown, ".pdf": parse_pdf, ".docx": parse_docx, ".pptx": parse_pptx}
    if suffix not in parsers:
        raise ValueError(f"不支持文本提取的类型: {suffix}")
    return clean_text(parsers[suffix](file_path))
